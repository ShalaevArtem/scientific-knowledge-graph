# -*- coding: utf-8 -*-
"""Презентация кейса «Научный клубок» — 10 слайдов, 16:9, python-pptx."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# палитра: графит + медь (горно-металлургическая тема)
INK = RGBColor(0x2A, 0x2D, 0x33)      # графит — тёмный фон/текст
COPPER = RGBColor(0xC0, 0x62, 0x2B)   # медь — акцент
SLATE = RGBColor(0x5A, 0x64, 0x72)    # приглушённый серо-синий
LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xF6, 0xEC, 0xE3)     # тёплая медная подложка карточек
TINT2 = RGBColor(0xEE, 0xF1, 0xF5)    # холодная подложка
GOOD = RGBColor(0x2C, 0x5F, 0x2D)
MUTED = RGBColor(0x8A, 0x92, 0x9E)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

HEAD = "Cambria"
BODY = "Calibri"


def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, x, y, w, h, fill=None, line=None, radius=None):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font=BODY, anchor=MSO_ANCHOR.TOP, space_after=6, line_spacing=None):
    """runs: строка или список абзацев; абзац = строка или (текст, {опции})."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        if isinstance(para, tuple):
            t, opt = para
        else:
            t, opt = para, {}
        pieces = t if isinstance(t, list) else [(t, {})]
        for j, (txt, ropt) in enumerate(pieces):
            r = p.add_run()
            r.text = txt
            f = r.font
            f.size = Pt(ropt.get("size", opt.get("size", size)))
            f.bold = ropt.get("bold", opt.get("bold", bold))
            f.italic = ropt.get("italic", opt.get("italic", False))
            f.color.rgb = ropt.get("color", opt.get("color", color))
            f.name = ropt.get("font", opt.get("font", font))
    return tb


def chip(s, x, y, w, h, label, fill=TINT2, color=INK, size=12, bold=False):
    b = box(s, x, y, w, h, fill=fill, radius=0.5)
    tf = b.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = BODY
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return b


def title(s, t, sub=None, color=INK):
    text(s, 0.6, 0.35, 12.1, 0.8, t, size=32, bold=True, font=HEAD, color=color)
    if sub:
        text(s, 0.6, 1.02, 12.1, 0.4, sub, size=14, color=SLATE)


def ball(s, x, y, d, fill, label=None, lsize=13, lcolor=LIGHT):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    c.shadow.inherit = False
    if label is not None:
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(lsize)
        r.font.bold = True
        r.font.color.rgb = lcolor
        r.font.name = BODY
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return c


def edge(s, x1, y1, x2, y2, color=SLATE, wpt=1.25):
    ln = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(wpt)
    ln.shadow.inherit = False
    return ln


def arrow(s, x, y, w=0.35, color=COPPER):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.22))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


# ============ 1. Титул ============
s = slide(INK)
# стилизованный «клубок» — маленький граф справа
nodes = [(10.4, 1.7, 0.55), (11.6, 1.15, 0.4), (12.15, 2.3, 0.5), (10.9, 3.05, 0.42),
         (11.95, 3.75, 0.55), (10.15, 4.35, 0.4), (12.6, 4.7, 0.36), (11.1, 5.35, 0.46)]
pairs = [(0, 1), (0, 2), (0, 3), (1, 2), (2, 3), (2, 4), (3, 4), (3, 5), (4, 6), (4, 7), (5, 7), (6, 7)]
for a, b in pairs:
    ax, ay, ad = nodes[a]
    bx, by, bd = nodes[b]
    edge(s, ax + ad / 2, ay + ad / 2, bx + bd / 2, by + bd / 2, color=RGBColor(0x55, 0x5C, 0x66))
for i, (x, y, d) in enumerate(nodes):
    ball(s, x, y, d, COPPER if i % 3 == 0 else RGBColor(0x9A, 0xA3, 0xAF))

text(s, 0.7, 2.0, 8.6, 1.4, "Научный клубок", size=60, bold=True, font=HEAD, color=LIGHT)
text(s, 0.7, 3.25, 8.6, 0.6, "Карта знаний R&D · горно-металлургические исследования",
     size=22, color=RGBColor(0xE8, 0xD5, 0xC4))
text(s, 0.7, 4.05, 8.8, 0.9,
     "GraphRAG: граф знаний (Neo4j) + гибридный семантический поиск + LLM-слой с цитатами",
     size=16, color=RGBColor(0xB9, 0xC0, 0xCA))
text(s, 0.7, 6.6, 8, 0.4, "Кейс «Единая карта знаний R&D» · июль 2026", size=13, color=MUTED)

# ============ 2. Проблема ============
s = slide()
title(s, "Проблема: знания R&D рассеяны",
      "институциональная память живёт в отчётах, презентациях и личных архивах — и не переиспользуется")
pains = [
    ("Потеря памяти", "методы обессоливания, циркуляция католита, распределение драгметаллов — в сотнях несвязанных файлов"),
    ("Дублирование", "команды заново делают литобзоры по очистке шахтных вод и удалению SO₂"),
    ("Нет связей", "не найти связь «кучное выщелачивание в холодном климате» → «выход металла»"),
    ("Медленные решения", "ответ про мировую практику подачи электролита = недели ручного поиска"),
    ("Противоречия", "нет верифицированной базы — конфликты интерпретаций (скорость циркуляции католита)"),
    ("Носители экспертизы", "неясно, кто в организации уже работал с аналогичной задачей"),
]
for i, (h, b) in enumerate(pains):
    cx, cy = 0.6 + (i % 3) * 4.15, 1.7 + (i // 3) * 2.55
    box(s, cx, cy, 3.95, 2.3, fill=TINT2 if i % 2 else TINT, radius=0.08)
    ball(s, cx + 0.25, cy + 0.28, 0.5, COPPER, label=str(i + 1), lsize=16)
    text(s, cx + 0.95, cy + 0.3, 2.8, 0.5, h, size=17, bold=True, font=HEAD)
    text(s, cx + 0.25, cy + 0.95, 3.45, 1.25, b, size=13, color=SLATE, line_spacing=1.05)

# ============ 3. Решение ============
s = slide()
title(s, "Решение: GraphRAG — граф знаний + поиск + LLM",
      "три слоя над единым хранилищем Neo4j")
cols = [
    ("Граф знаний", "скелет связей",
     ["13 типов узлов: материалы, процессы, оборудование, предприятия, измерения, выводы, эксперты",
      "provenance каждого факта: документ, страница/слайд, модель, уверенность",
      "RU/EN-синонимы: «ПВП» ≡ «flash furnace»"]),
    ("Гибридный поиск", "4 канала, слияние RRF",
     ["вектор (эмбеддинги RU/EN) — перефразировки; BM25 — марки и формулы: «12Х18Н10Т», «SO₂»",
      "граф: документы, покрывающие несколько сущностей вопроса (рёбра MENTIONS)",
      "измерения: числовые ограничения вопроса → интервалы значений, «мг/дм³» ≡ «мг/л»"]),
    ("LLM-слой", "вопрос → ответ с цитатами",
     ["разбор вопроса в фильтры («за 5 лет» → 2021–2026)",
      "сводка, где каждое утверждение подкреплено [n]-цитатой",
      "выделение расхождений источников и пробелов"]),
]
for i, (h, sub, items) in enumerate(cols):
    cx = 0.6 + i * 4.3
    box(s, cx, 1.65, 4.05, 3.6, fill=TINT2, radius=0.06)
    text(s, cx + 0.3, 1.95, 3.5, 0.5, h, size=20, bold=True, font=HEAD, color=COPPER)
    text(s, cx + 0.3, 2.42, 3.5, 0.35, sub, size=13, color=SLATE, space_after=2)
    text(s, cx + 0.3, 2.85, 3.5, 2.3,
         [(t, {}) for t in items], size=12.5, color=INK, space_after=8, line_spacing=1.02)
box(s, 0.6, 5.55, 6.1, 1.35, fill=TINT, radius=0.08)
text(s, 0.9, 5.75, 5.6, 1.0,
     [([("Инвариант 1 · ", {"bold": True, "color": COPPER}),
        ("«нет источника — нет факта»", {"bold": True})], {}),
      ("каждый факт в графе ссылается на документ и страницу", {"color": SLATE, "size": 12.5})],
     size=14)
box(s, 6.95, 5.55, 5.75, 1.35, fill=TINT, radius=0.08)
text(s, 7.25, 5.75, 5.3, 1.0,
     [([("Инвариант 2 · ", {"bold": True, "color": COPPER}),
        ("верификация", {"bold": True})], {}),
      ("новые сущности и выводы из текстов идут в очередь эксперта (needs_review)", {"color": SLATE, "size": 12.5})],
     size=14)

# ============ 4. Онтология ============
s = slide()
title(s, "Онтология предметной области", "13 типов узлов, ключевые связи — с provenance на рёбрах")
labels = ["Material", "Process", "Equipment", "Facility", "Document", "Chunk", "Property",
          "Measurement", "Conclusion", "Experiment", "Regime", "Person", "TopicTag"]
for i, lab in enumerate(labels):
    chip(s, 0.6 + (i % 4) * 1.62, 1.75 + (i // 4) * 0.62, 1.5, 0.46, lab,
         fill=TINT if i % 2 else TINT2, size=12, bold=True)
text(s, 0.6, 4.55, 6.0, 2.4,
     [("Ключевые связи:", {"bold": True}),
      ("Document —MENTIONS→ Material / Process / Equipment / Facility", {"size": 12.5, "color": SLATE}),
      ("Document —REPORTS→ Measurement —OF_PROPERTY→ Property", {"size": 12.5, "color": SLATE}),
      ("Conclusion —DESCRIBED_IN→ Document, —ABOUT→ сущность", {"size": 12.5, "color": SLATE}),
      ("Document —AUTHORED_BY→ Person, —TAGGED→ TopicTag", {"size": 12.5, "color": SLATE}),
      ("Measurement.conditions: параметр / оператор / диапазон («сульфаты ≤ 300 мг/л»)", {"size": 12.5, "color": SLATE})],
     size=14, space_after=7)
# мини-схема справа
gx = 7.4
gn = {
    "Document": (gx + 1.9, 1.8, 1.5, INK),
    "Material": (gx + 0.1, 3.1, 1.4, COPPER),
    "Process": (gx + 2.0, 3.6, 1.35, COPPER),
    "Measurement": (gx + 3.9, 3.1, 1.6, SLATE),
    "Conclusion": (gx + 0.6, 4.9, 1.5, SLATE),
    "Person": (gx + 4.0, 1.9, 1.2, SLATE),
    "Property": (gx + 4.1, 4.6, 1.35, MUTED),
}
rels = [("Document", "Material", "MENTIONS"), ("Document", "Process", "MENTIONS"),
        ("Document", "Measurement", "REPORTS"), ("Measurement", "Property", "OF_PROPERTY"),
        ("Conclusion", "Document", "DESCRIBED_IN"), ("Conclusion", "Material", "ABOUT"),
        ("Document", "Person", "AUTHORED_BY")]
cpt = {k: (x + w / 2, y + 0.27) for k, (x, y, w, c) in gn.items()}
for a, b, _ in rels:
    edge(s, cpt[a][0], cpt[a][1], cpt[b][0], cpt[b][1], color=MUTED, wpt=1.0)
for k, (x, y, w, c) in gn.items():
    chip(s, x, y, w, 0.54, k, fill=c, color=LIGHT, size=12, bold=True)

# ============ 5. Конвейер ============
s = slide()
title(s, "Конвейер: от папки с документами до ответа", "все шаги возобновляемы; LLM-провайдер подключаемый")
steps = [
    ("Корпус", "PDF · DOCX · PPTX\n4,5 ГБ, 1 522 документа,\nRU/EN, 2001–2026"),
    ("skg ingest", "текст с привязкой\nк страницам/слайдам →\nDocument + Chunk"),
    ("skg embed", "локальные эмбеддинги\n(fastembed ONNX, RU/EN),\nвекторный индекс + BM25"),
    ("skg extract", "LLM по онтологии\n(Yandex AI Studio),\nнормализация, needs_review"),
    ("Слой запросов", "вопрос → фильтры →\nгибридный поиск →\nсводка с цитатами"),
    ("Веб-UI", "ответ и источники,\nграф, пробелы,\nэкспорт Markdown"),
]
for i, (h, b) in enumerate(steps):
    cx = 0.55 + i * 2.13
    box(s, cx, 2.3, 1.95, 2.5, fill=TINT if i % 2 else TINT2, radius=0.09)
    text(s, cx + 0.16, 2.5, 1.65, 0.6, h, size=15, bold=True, font=HEAD, color=COPPER,
         align=PP_ALIGN.CENTER)
    text(s, cx + 0.14, 3.15, 1.68, 1.5, b.replace("\n", " "), size=11.5, color=INK,
         align=PP_ALIGN.CENTER, line_spacing=1.05)
    if i < 5:
        arrow(s, cx + 1.95, 3.42, w=0.19)
text(s, 0.6, 5.35, 12.1, 0.5,
     [([("Хранилище — Neo4j: ", {"bold": True}),
        ("граф, векторный индекс и полнотекстовый поиск в одной БД — без отдельного поискового движка",
         {"color": SLATE})], {})], size=14)
text(s, 0.6, 6.0, 12.1, 0.9,
     [([("Отказоустойчивость: ", {"bold": True}),
        ("битые PDF и сканы не валят импорт (5 из 1 527 отклонено с диагностикой); ошибки LLM ретраятся, "
         "невалидный JSON — повторный запрос; упавшие чанки добираются повторным прогоном",
         {"color": SLATE})], {})], size=14)

# ============ 6. Демо: ответ с цитатами ============
s = slide()
title(s, "Демо · вопрос из ТЗ → ответ с цитатами",
      "«какие решения циркуляции католита при электроэкстракции никеля описаны в мировой практике?»")
box(s, 0.6, 1.75, 8.1, 5.1, fill=TINT2, radius=0.05)
text(s, 0.9, 2.0, 7.5, 4.6,
     [("Применяются: разделение катодного и анодного пространства диафрагмой и регулирование "
       "подачи католита для поддержания перепада уровня между католитом и анолитом.", {"size": 13}),
      ([("С ростом скорости от 0,5 до 4,9 л/мин шероховатость катодной меди снижается "
         "с 6,86 до 2,97 мкм ", {"size": 13}),
        ("[2]", {"size": 13, "bold": True, "color": COPPER})], {}),
      ([("Контроль скорости — для положительного гидростатического напора в катодном мешке "
         "и снижения массопереноса ионов водорода ", {"size": 13}),
        ("[4]", {"size": 13, "bold": True, "color": COPPER})], {}),
      ([("При плотности тока 250 А/м² скорость циркуляции — 25 дм³/А·ч ", {"size": 13}),
        ("[5]", {"size": 13, "bold": True, "color": COPPER}),
        (";  типовая скорость 20–30 л/ч ", {"size": 13}),
        ("[6]", {"size": 13, "bold": True, "color": COPPER})], {}),
      ([("Пробелы/что уточнить: ", {"size": 13, "bold": True}),
        ("влияние скорости потока на эффективность электроэкстракции освещено слабо",
         {"size": 13, "color": SLATE})], {})],
     space_after=10, line_spacing=1.1)
feats = [("Числа — точно из источника", "концентрации и скорости не «сочиняются»"),
         ("Клик по [n] — фрагмент и контекст", "±2 соседних чанка, путь к файлу"),
         ("Фильтры из вопроса", "«за последние 5 лет» → 2021–2026"),
         ("Экспорт в Markdown", "готовый фрагмент для отчёта/ТЗ")]
for i, (h, b) in enumerate(feats):
    cy = 1.75 + i * 1.32
    box(s, 8.95, cy, 3.8, 1.15, fill=TINT, radius=0.09)
    text(s, 9.2, cy + 0.14, 3.35, 0.5, h, size=13.5, bold=True)
    text(s, 9.2, cy + 0.62, 3.35, 0.45, b, size=11.5, color=SLATE)

# ============ 7. Демо: граф и пробелы ============
s = slide()
title(s, "Демо · граф, пробелы, очередь эксперта",
      "окрестность сущности · тепловая карта «материал × процесс» · needs_review")
# слева — мини-граф
box(s, 0.6, 1.7, 5.9, 5.2, fill=TINT2, radius=0.05)
text(s, 0.9, 1.92, 5.3, 0.4, "Граф: «электроэкстракция» — 92 узла, 135 связей",
     size=14, bold=True)
gn2 = [("электроэкстракция", 2.7, 3.6, 2.0, COPPER), ("никель", 1.0, 2.7, 1.2, INK),
       ("католит", 1.15, 4.9, 1.2, INK), ("диафрагменная ячейка", 4.2, 2.55, 1.9, SLATE),
       ("Nikkelverk", 4.6, 5.3, 1.3, SLATE), ("медь", 0.95, 3.8, 1.0, SLATE),
       ("Cvet. metally 2010", 3.0, 5.75, 1.7, MUTED)]
c2 = {n: (x + w / 2, y + 0.25) for n, x, y, w, c in gn2}
for a, b in [("электроэкстракция", "никель"), ("электроэкстракция", "католит"),
             ("электроэкстракция", "диафрагменная ячейка"), ("электроэкстракция", "Nikkelverk"),
             ("электроэкстракция", "медь"), ("электроэкстракция", "Cvet. metally 2010")]:
    edge(s, c2[a][0], c2[a][1], c2[b][0], c2[b][1], color=MUTED, wpt=1.0)
for n, x, y, w, c in gn2:
    chip(s, x, y, w, 0.5, n, fill=c, color=LIGHT, size=11, bold=True)
# справа — тепловая карта
box(s, 6.85, 1.7, 5.9, 3.5, fill=TINT2, radius=0.05)
text(s, 7.15, 1.92, 5.3, 0.4, "Пробелы: документов на пару «материал × процесс»",
     size=14, bold=True)
mats = ["никель", "медь", "золото", "гипс"]
procs = ["выщелач.", "электроэкстр.", "флотация", "обжиг"]
vals = [[14, 15, 6, 9], [11, 3, 8, 7], [9, 0, 5, 2], [4, 0, 0, 1]]
for j, pname in enumerate(procs):
    text(s, 8.6 + j * 1.02, 2.35, 1.0, 0.3, pname, size=9.5, color=SLATE, align=PP_ALIGN.CENTER)
for i, m in enumerate(mats):
    text(s, 7.15, 2.72 + i * 0.56, 1.4, 0.3, m, size=10.5, color=SLATE)
    for j in range(4):
        v = vals[i][j]
        f = RGBColor(0xD9, 0x53, 0x3C) if v == 0 else (
            RGBColor(0xC0 - min(v, 15) * 4, 0x62 + min(v, 15) * 3, 0x2B + min(v, 15) * 4))
        cell = box(s, 8.62 + j * 1.02, 2.68 + i * 0.56, 0.92, 0.48,
                   fill=(TINT if v else RGBColor(0xF3, 0xC9, 0xC0)), radius=0.15)
        tfc = cell.text_frame
        tfc.margin_left = tfc.margin_right = tfc.margin_top = tfc.margin_bottom = 0
        pc = tfc.paragraphs[0]
        pc.alignment = PP_ALIGN.CENTER
        rc = pc.add_run()
        rc.text = str(v)
        rc.font.size = Pt(12)
        rc.font.bold = (v == 0)
        rc.font.color.rgb = RGBColor(0xB0, 0x28, 0x18) if v == 0 else INK
        tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
text(s, 7.15, 4.9, 5.4, 0.3, "красные нули — неизученные комбинации → кандидаты на НИР",
     size=11, color=SLATE)
box(s, 6.85, 5.35, 5.9, 1.55, fill=TINT, radius=0.08)
text(s, 7.15, 5.55, 5.4, 1.2,
     [([("Очередь эксперта: ", {"bold": True}),
        ("все сущности и выводы, впервые извлечённые из текстов, помечены needs_review",
         {"color": SLATE, "size": 12.5})], {}),
      ("эксперт подтверждает или правит — дата и автор фиксируются", {"color": SLATE, "size": 12.5})],
     size=13.5)

# ============ 8. Цифры ============
s = slide(INK)
title(s, "Система в цифрах", None, color=LIGHT)
text(s, 0.6, 1.02, 12, 0.4, "реальный корпус кейса, обработан локально + Yandex AI Studio",
     size=14, color=MUTED)
stats = [
    ("1 522", "документа импортировано", "PDF/DOCX/PPTX · 5 категорий · 2001–2026 · RU/EN"),
    ("290 508", "фрагментов в графе", "с привязкой к странице/слайду; 100 % с эмбеддингами"),
    ("74 100+", "измерений с условиями", "значение · единица · диапазон · условия процесса"),
    ("15 300+", "выводов с географией", "РФ/зарубежная практика, уверенность, источник"),
    ("42 000+", "сущностей", "материалы · процессы · оборудование · предприятия · эксперты"),
    ("~0,5 c", "поиск, 4 канала", "вектор + BM25 + граф + измерения; сводка LLM ~20–40 с"),
]
for i, (num, lab, sub) in enumerate(stats):
    cx, cy = 0.6 + (i % 3) * 4.25, 1.75 + (i // 3) * 2.6
    box(s, cx, cy, 4.0, 2.3, fill=RGBColor(0x35, 0x39, 0x41), radius=0.08)
    text(s, cx + 0.3, cy + 0.25, 3.4, 0.9, num, size=40, bold=True, font=HEAD, color=COPPER)
    text(s, cx + 0.3, cy + 1.15, 3.4, 0.4, lab, size=15, bold=True, color=LIGHT)
    text(s, cx + 0.3, cy + 1.6, 3.4, 0.6, sub, size=11.5, color=RGBColor(0xB9, 0xC0, 0xCA))

# ============ 9. Соответствие требованиям ============
s = slide()
title(s, "Соответствие требованиям кейса")
rows = [
    ("Многопараметрические запросы", "материал + процесс + условия + география + период; графовый канал поиска (MENTIONS)", True),
    ("Верификация и версии знаний", "provenance до страницы, уверенность, очередь needs_review; история изменений выводов", True),
    ("Отечественная / зарубежная практика", "география у выводов и предприятий, разделение в ответах LLM", True),
    ("Числовые ограничения и диапазоны", "извлечение и поиск: интервалы значений, нормализация единиц («мг/дм³» ≡ «мг/л»)", True),
    ("Масштабируемость доменов", "онтология расширяется без миграций; новые категории = папка корпуса", True),
    ("Мультиязычность RU/EN", "мультиязычные эмбеддинги, синонимы терминов, «ПВП» ≡ «flash furnace»", True),
    ("RBAC и аудит действий", "5 ролей; внешний партнёр не видит внутренние категории; журнал запросов и просмотров", True),
    ("Экспорт, уведомления", "Markdown в UI; PDF/JSON-LD и подписки на темы — в road map", None),
]
for i, (h, b, ok) in enumerate(rows):
    cy = 1.55 + i * 0.68
    mark = "✓" if ok else "→"
    ball(s, 0.6, cy + 0.04, 0.42, GOOD if ok else SLATE, label=mark, lsize=15)
    text(s, 1.25, cy, 4.35, 0.6, h, size=14.5, bold=True)
    text(s, 5.75, cy + 0.03, 7.0, 0.6, b, size=12.5, color=SLATE)

# ============ 10. Roadmap и ограничения ============
s = slide()
title(s, "Честные ограничения и road map")
box(s, 0.6, 1.6, 6.0, 5.3, fill=TINT2, radius=0.06)
text(s, 0.9, 1.85, 5.4, 0.5, "Ограничения сегодня", size=19, bold=True, font=HEAD, color=SLATE)
text(s, 0.9, 2.5, 5.4, 4.2,
     [("качество извлечения = качество LLM: ~0,4 % чанков модель не разобрала даже с ретраями", {}),
      ("сканы без текстового слоя (5 документов) требуют OCR — не подключён", {}),
      ("извлечение сущностей прошло по приоритетным категориям; журнальные подшивки — по мере бюджета API", {}),
      ("SUPPORTS/CONTRADICTS между выводами из текстов размечаются пока вручную", {}),
      ("доступ разграничен по категориям документов; гриф на уровне отдельного документа — в road map", {})],
     size=13.5, color=INK, space_after=10, line_spacing=1.05)
box(s, 6.85, 1.6, 5.9, 5.3, fill=TINT, radius=0.06)
text(s, 7.15, 1.85, 5.3, 0.5, "Road map", size=19, bold=True, font=HEAD, color=COPPER)
text(s, 7.15, 2.5, 5.3, 4.2,
     [("автоматическая разметка противоречий (CONTRADICTS) и консенсуса", {}),
      ("гриф доступа на уровне отдельного документа, SSO вместо токенов", {}),
      ("экспорт PDF и JSON-LD, подписки-уведомления на новые публикации по теме", {}),
      ("дашборды покрытия по направлениям: гидрометаллургия, экология, отходы", {}),
      ("OCR-слой для сканов, данные с датчиков установок как новый тип источника", {})],
     size=13.5, color=INK, space_after=10, line_spacing=1.05)

out = sys.argv[1] if len(sys.argv) > 1 else "presentation.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
