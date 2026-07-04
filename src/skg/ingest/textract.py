"""Извлечение текста из PDF / DOCX / DOCM / PPTX / XLS / XLSX / XLSM.

Возвращает список юнитов (страница, слайд, блок абзацев или лист таблицы) с их
текстом — номер юнита попадает в provenance чанка, чтобы цитату можно было найти
в оригинале.
"""

from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif"}
SUPPORTED = {".pdf", ".docx", ".docm", ".pptx", ".xls", ".xlsx", ".xlsm"} | IMAGE_EXTS

# защита от патологически больших книг: строк на лист и символов на лист
MAX_ROWS_PER_SHEET = 5000
MAX_CHARS_PER_SHEET = 200_000

# OCR (Tesseract) — опционально, включается SKG_OCR=1. Медленный, поэтому по
# умолчанию выключен: обычный ingest не тормозит. Когда включён — распознаёт
# страницы PDF без текстового слоя (сканы, картиночные вставки) и файлы-изображения.
OCR_LANG = os.environ.get("SKG_OCR_LANG", "rus+eng")
OCR_DPI = int(os.environ.get("SKG_OCR_DPI", "300"))
# минимум символов на странице, ниже которого считаем текстовый слой отсутствующим
OCR_MIN_CHARS = int(os.environ.get("SKG_OCR_MIN_CHARS", "20"))


@dataclass
class TextUnit:
    unit_no: int  # страница PDF, номер слайда, порядковый блок DOCX, лист таблицы
    unit_kind: str  # "page" | "slide" | "block" | "sheet" | "image"
    text: str


class ExtractionError(Exception):
    pass


def ocr_enabled() -> bool:
    return os.environ.get("SKG_OCR", "").lower() in ("1", "true", "on", "yes")


_ocr_ready: bool | None = None


def _configure_tesseract() -> None:
    """Автонастройка без правки PATH/системного tessdata (нет прав на Program Files):
    бинарь ищем по TESSERACT_CMD или стандартному пути Windows; языковые данные —
    в локальном .cache/tessdata (rus+eng+osd лежат там), если TESSDATA_PREFIX не задан."""
    import pytesseract

    cmd = os.environ.get("TESSERACT_CMD")
    if not cmd:
        default = Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe")
        if default.exists():
            cmd = str(default)
    if cmd:
        pytesseract.pytesseract.tesseract_cmd = cmd

    if not os.environ.get("TESSDATA_PREFIX"):
        local = Path(__file__).resolve().parents[3] / ".cache" / "tessdata"
        if (local / "rus.traineddata").exists():
            os.environ["TESSDATA_PREFIX"] = str(local)


def _ocr_ready_once() -> bool:
    """Проверяет доступность Tesseract один раз за процесс. Если OCR включён, но
    бинаря/пакета нет — предупреждаем и работаем без OCR, не роняя пайплайн."""
    global _ocr_ready
    if _ocr_ready is None:
        try:
            import pytesseract

            _configure_tesseract()
            pytesseract.get_tesseract_version()
            _ocr_ready = True
        except Exception as exc:
            print(
                f"[OCR] SKG_OCR включён, но Tesseract недоступен ({exc}); "
                "продолжаю без OCR. Установите бинарь Tesseract и языковые данные rus+eng.",
                flush=True,
            )
            _ocr_ready = False
    return _ocr_ready


def _ocr_image(image) -> str:
    import pytesseract

    return pytesseract.image_to_string(image, lang=OCR_LANG).strip()


def extract_text(path: Path) -> list[TextUnit]:
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return _extract_pdf(path)
        if suffix in (".docx", ".docm"):
            return _extract_docx(path)
        if suffix == ".pptx":
            return _extract_pptx(path)
        if suffix in (".xlsx", ".xlsm"):
            return _extract_xlsx(path)
        if suffix == ".xls":
            return _extract_xls(path)
        if suffix in IMAGE_EXTS:
            return _extract_image(path)
    except Exception as exc:  # повреждённые файлы в корпусе не должны валить пайплайн
        raise ExtractionError(f"{path.name}: {exc}") from exc
    raise ExtractionError(f"{path.name}: неподдерживаемый формат {suffix}")


def _extract_pdf(path: Path) -> list[TextUnit]:
    import fitz  # pymupdf

    ocr = ocr_enabled() and _ocr_ready_once()
    units = []
    with fitz.open(path) as doc:
        for page_no, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            # текстового слоя нет (скан/картиночная страница) — пробуем OCR
            if ocr and len(text) < OCR_MIN_CHARS and page.get_images():
                text = _ocr_pdf_page(page) or text
            if text:
                units.append(TextUnit(page_no, "page", text))
    return units


def _ocr_pdf_page(page) -> str:
    """Рендерит страницу PDF в растр и распознаёт Tesseract-ом. Ошибки OCR не
    должны валить извлечение всего документа — при сбое возвращаем пустую строку."""
    from io import BytesIO

    from PIL import Image

    try:
        pix = page.get_pixmap(dpi=OCR_DPI)
        with Image.open(BytesIO(pix.tobytes("png"))) as img:
            return _ocr_image(img)
    except Exception as exc:
        print(f"[OCR] страница не распознана: {exc}", flush=True)
        return ""


def _extract_image(path: Path) -> list[TextUnit]:
    """Файл-изображение (скан/фотография страницы) → текст через OCR.
    Без включённого OCR изображения пропускаются (текста не извлечь)."""
    if not (ocr_enabled() and _ocr_ready_once()):
        return []
    from PIL import Image

    with Image.open(path) as img:
        text = _ocr_image(img)
    return [TextUnit(1, "image", text)] if text else []


def _open_word(path: Path):
    """Открывает docx/docm. python-docx отказывает macro-enabled документам по
    content-type — подменяем тип в [Content_Types].xml на обычный (макросы не нужны,
    текст тот же)."""
    import docx

    try:
        return docx.Document(str(path))
    except ValueError:
        import io
        import zipfile

        src = zipfile.ZipFile(path)
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as dst:
            for item in src.infolist():
                data = src.read(item)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"application/vnd.ms-word.document.macroEnabled.main+xml",
                        b"application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
                    )
                dst.writestr(item, data)
        buf.seek(0)
        return docx.Document(buf)


def _extract_docx(path: Path, block_paragraphs: int = 40) -> list[TextUnit]:
    doc = _open_word(path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # таблицы тоже несут данные (составы, параметры) — добавляем построчно
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    units = []
    for i in range(0, len(paragraphs), block_paragraphs):
        block = "\n".join(paragraphs[i : i + block_paragraphs])
        units.append(TextUnit(i // block_paragraphs + 1, "block", block))
    return units


def _extract_pptx(path: Path) -> list[TextUnit]:
    from pptx import Presentation

    prs = Presentation(str(path))
    units = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            units.append(TextUnit(slide_no, "slide", "\n".join(texts)))
    return units


def _cell_str(value) -> str:
    """Ячейка → строка. Целые float («1.0») печатаем как «1», чтобы не мусорить
    числовые данные; None/пустые — пустой строкой (отбрасываются вызывающим)."""
    if value is None:
        return ""
    if isinstance(value, bool):
        return "да" if value else "нет"
    if isinstance(value, float):
        return f"{value:g}"
    return str(value).strip()


def _rows_to_unit(sheet_no: int, sheet_name: str, rows: list[list[str]]) -> TextUnit | None:
    """Строки листа → один юнит (unit_kind="sheet"). Пустые ячейки и строки
    отбрасываются; заголовок с именем листа делает лист находимым по названию."""
    lines: list[str] = []
    chars = 0
    for cells in rows:
        clean = [c for c in cells if c]
        if not clean:
            continue
        line = " | ".join(clean)
        lines.append(line)
        chars += len(line) + 1
        if chars >= MAX_CHARS_PER_SHEET:
            lines.append("… (лист усечён)")
            break
    if not lines:
        return None
    header = f"Лист: {sheet_name}" if sheet_name else f"Лист {sheet_no}"
    return TextUnit(sheet_no, "sheet", header + "\n" + "\n".join(lines))


def _extract_xlsx(path: Path) -> list[TextUnit]:
    from openpyxl import load_workbook

    # read_only — не держим книгу в памяти; data_only — берём вычисленные значения
    # формул, а не сами формулы (без пересчёта показал бы «=A1*B1»).
    wb = load_workbook(str(path), read_only=True, data_only=True)
    try:
        units = []
        for sheet_no, ws in enumerate(wb.worksheets, start=1):
            rows = []
            for r, row in enumerate(ws.iter_rows(values_only=True)):
                if r >= MAX_ROWS_PER_SHEET:
                    break
                rows.append([_cell_str(v) for v in row])
            unit = _rows_to_unit(sheet_no, ws.title, rows)
            if unit:
                units.append(unit)
        return units
    finally:
        wb.close()


def _extract_xls(path: Path) -> list[TextUnit]:
    import xlrd

    book = xlrd.open_workbook(str(path))
    units = []
    for sheet_no, sh in enumerate(book.sheets(), start=1):
        rows = []
        for r in range(min(sh.nrows, MAX_ROWS_PER_SHEET)):
            rows.append([_cell_str(v) for v in sh.row_values(r)])
        unit = _rows_to_unit(sheet_no, sh.name, rows)
        if unit:
            units.append(unit)
    return units
